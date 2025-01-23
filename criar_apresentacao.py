from pptx.enum.text import PP_ALIGN
from pptx import Presentation
from pptx.util import Inches
from datetime import datetime


def criar_apresentacao(topicos_separados,caminho_arquivo, tema):
# Carrega o arquivo de template
    data_atual = datetime.now().strftime('%d/%m/%Y')
    arquivo = 'apresentacao.pptx'
    apresentacao = Presentation('modelo_octopus.pptx')

    for titulo, subticos in topicos_separados.items():
        # Cria um slide separando para cada título
        slide_layout = apresentacao.slide_layouts[1]  # Escolhe um layout (aqui o de título e conteúdo)
        slide = apresentacao.slides.add_slide(slide_layout)
        
        # Configura o título do slide
        titulo_shape = slide.shapes.title
        titulo_shape.text = titulo

        # Define o texto do título como negrito
        for run in titulo_shape.text_frame.paragraphs[0].runs:
            run.bold = True  # Define o texto em negrito

        titulo_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT  # Alinha à esquerda
        
        # Para cada subtítulo, adiciona os conteúdos com marcadores
        for subtitulo, conteudo in subticos.items():
            # Adiciona o subtítulo como um novo parágrafo
            caixa_texto = slide.shapes.placeholders[1]  # Placeholder 1 (geralmente é a caixa de conteúdo)

            # Adiciona o subtítulo como um novo parágrafo
            p_subtitulo = caixa_texto.text_frame.add_paragraph()  
            p_subtitulo.text = subtitulo
            p_subtitulo.space_after = Inches(0.1)  # Espaçamento após o subtítulo

            # Define o texto do subtítulo como negrito
            for run in p_subtitulo.runs:
                run.font.bold = True  # Define o texto do subtítulo em negrito
            
            # Adiciona o conteúdo como uma lista (com marcadores)
            for item in conteudo.split('\n'):
                if item.strip():  # Verifica se a linha não está vazia
                    p_conteudo = caixa_texto.text_frame.add_paragraph()  # Adiciona um novo parágrafo para cada item
                    p_conteudo.text = item.strip()  # Adiciona o texto do conteúdo
                    p_conteudo.space_after = Inches(0.05)  # Espaçamento após cada item
                    p_conteudo.level = 0  # Define o nível do marcador (0 é o nível principal)
                    p_conteudo.bullet = True  # Adiciona marcador ao parágrafo

        # Ajusta o alinhamento e formatação do caixa de texto
        for par in caixa_texto.text_frame.paragraphs:
            par.alignment = PP_ALIGN.LEFT 
            for run in par.runs:
                run.font.size = Inches(0.4) 
    salva_arquivo = caminho_arquivo+'/apresentacao_'+tema.replace(' ','-')+'.pptx'
    nome_arquivo = 'apresentacao_'+tema.replace(' ', '-')+'.pptx'
    print(f'Nome arquivo: {nome_arquivo}')
    print(f'Caminho de salvamento do arquivo: {salva_arquivo}')
    # Salva a apresentação em um novo arquivo
    apresentacao.save(salva_arquivo)
    print('Apresentação criada com sucesso!')
    return nome_arquivo