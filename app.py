from pptx.util import Inches  
from pptx.enum.text import PP_ALIGN
from flask import Flask, request, send_file, render_template, jsonify
from pptx import Presentation
from openai import OpenAI
from separa_topico import *
from helpers import *

modelo = carrega('modelo_apresentacao.txt')
app = Flask(__name__)

#client = OpenAI(api_key='Chave API AQUI')

def gerar_conteudo_gpt(prompt):
    try:
        completion = client.chat.completions.create(
        model="gpt-4o",
        messages=[
            {"role": "system", "content": 
            f"Você atuará como uma ferramenta de criação de apresentação, voce recebera um prompt e a partir deste prompt voce vai criar uma apresentação de power point. A sua resposta será renderizada em um slide, então evite usar formatação markdown. Nas suas respostas não inclua indicativos dos slides como por exemplo: 'Slide 1: ,Slide 2:' etc. Utilize como exemplo para criar suas apresentações o seguinte modelo: {modelo}"},
            {
                "role": "user",
                "content": prompt
            }
        ]
    )
        # Retornar o conteúdo gerado
        #print(completion.choices[0].message.content)
        return completion.choices[0].message.content # O retorno desta função é String
    except Exception as e:
        print(f"Erro ao chamar a API: {e}")
        return None  # Retornando None em caso de erro
    

def criar_apresentacao(topicos_separados):
# Carrega o arquivo de template
    arquivo = 'apresentacao.pptx'
    apresentacao = Presentation('modelo_octopus.pptx')

    for titulo, subticos in topicos_separados.items():
        # Cria um slide separando para cada título
        slide_layout = apresentacao.slide_layouts[1]  # Escolhe um layout (aqui o de título e conteúdo)
        slide = apresentacao.slides.add_slide(slide_layout)
        
        # Configura o título do slide
        titulo_shape = slide.shapes.title
        titulo_shape.text = titulo

        # Define o texto do título como negrito
        for run in titulo_shape.text_frame.paragraphs[0].runs:
            run.bold = True  # Define o texto em negrito

        titulo_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT  # Alinha à esquerda
        
        # Para cada subtítulo, adiciona os conteúdos com marcadores
        for subtitulo, conteudo in subticos.items():
            # Adiciona o subtítulo como um novo parágrafo
            caixa_texto = slide.shapes.placeholders[1]  # Placeholder 1 (geralmente é a caixa de conteúdo)

            # Adiciona o subtítulo como um novo parágrafo
            p_subtitulo = caixa_texto.text_frame.add_paragraph()  
            p_subtitulo.text = subtitulo
            p_subtitulo.space_after = Inches(0.1)  # Espaçamento após o subtítulo

            # Define o texto do subtítulo como negrito
            for run in p_subtitulo.runs:
                run.font.bold = True  # Define o texto do subtítulo em negrito
            
            # Adiciona o conteúdo como uma lista (com marcadores)
            for item in conteudo.split('\n'):
                if item.strip():  # Verifica se a linha não está vazia
                    p_conteudo = caixa_texto.text_frame.add_paragraph()  # Adiciona um novo parágrafo para cada item
                    p_conteudo.text = item.strip()  # Adiciona o texto do conteúdo
                    p_conteudo.space_after = Inches(0.05)  # Espaçamento após cada item
                    p_conteudo.level = 0  # Define o nível do marcador (0 é o nível principal)
                    p_conteudo.bullet = True  # Adiciona marcador ao parágrafo

        # Ajusta o alinhamento e formatação do caixa de texto
        for par in caixa_texto.text_frame.paragraphs:
            par.alignment = PP_ALIGN.LEFT  # Alinha à esquerda
            for run in par.runs:
                run.font.size = Inches(0.4)  # Define o tamanho da fonte
                
    # Salva a apresentação em um novo arquivo
    apresentacao.save('apresentacao.pptx')
    print('Apresentação criada com sucesso!')
    return arquivo

# Rota para a página principal
@app.route('/')
def index():
    return render_template('index.html') 

@app.route('/frutas', methods=['GET'])
def get_frutas():
    frutas = ['Maçã', 'Banana', 'Laranja', 'Manga', 'Uva']
    app.logger.info("Rota de frutas acessada.")
    return jsonify(frutas)

# Rota para gerar a apresentação
@app.route('/criar', methods=['POST'])
def criar():
    tema = request.get_json['tema']
    texto = gerar_conteudo_gpt(prompt=tema)
    topicos = separa_topicos(texto)  # Separa os tópicos
    arquivo = criar_apresentacao(topicos)  # Cria a apresentação
    return send_file(arquivo, as_attachment=True, download_name='apresentacao.pptx', mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')  # Retorna o arquivo para download

if __name__ == '__main__':
    app.run(debug=True)
